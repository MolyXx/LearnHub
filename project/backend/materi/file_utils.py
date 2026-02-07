import os
import mimetypes
from pypdf import PdfReader
from pptx import Presentation
import docx
import openpyxl
import csv
import io
import traceback
from .helper import supabase_signed_to_public

# Ensure mimetypes are initialized
mimetypes.init()

def extract_file_content(file_field):
    """
    Extracts text and image URLs from a Django FileField.
    Returns: (text_content, image_urls)
    """
    text_content = ""
    image_urls = []
    
    if not file_field:
        return text_content, image_urls

    try:
        filename = file_field.name.lower()
        file_ext = os.path.splitext(filename)[1]
        
        # Determine mime type if possible, or fallback to extension
        mime_type, _ = mimetypes.guess_type(filename)
        
        # Special handling for images
        if mime_type and mime_type.startswith('image/'):
            # Check if using Supabase/S3 storage which provides a URL
            if hasattr(file_field, 'url'):
                # Convert to public URL if needed (for Supabase signed URLs)
                public_url = supabase_signed_to_public(file_field.url)
                if public_url:
                    image_urls.append(public_url)
            return f"[Gambar: {os.path.basename(filename)}]", image_urls

        # Open file in appropriate mode
        # Note: file_field.open() usually returns file-like object.
        # For text based files we might need specific encoding.
        
        if file_ext == ".pdf":
            with file_field.open('rb') as f:
                reader = PdfReader(f)
                pdf_text = []
                for i, page in enumerate(reader.pages):
                    page_text = page.extract_text()
                    if page_text:
                        pdf_text.append(f"Page {i+1}:\n{page_text}")
                text_content = "\n".join(pdf_text)

        elif file_ext == ".docx":
            with file_field.open('rb') as f:
                doc = docx.Document(f)
                text_content = "\n".join([para.text for para in doc.paragraphs])

        elif file_ext == ".pptx":
            with file_field.open('rb') as f:
                prs = Presentation(f)
                slides_text = []
                for i, slide in enumerate(prs.slides):
                    slide_text = []
                    # Extract title
                    if slide.shapes.title:
                        slide_text.append(f"Title: {slide.shapes.title.text}")
                    # Extract text from shapes
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text:
                             # Avoid duplicating title if it's already extracted
                            if shape == slide.shapes.title:
                                continue
                            slide_text.append(shape.text)
                    slides_text.append(f"Slide {i+1}:\n" + "\n".join(slide_text))
                text_content = "\n".join(slides_text)

        elif file_ext == ".xlsx":
            with file_field.open('rb') as f:
                wb = openpyxl.load_workbook(f, data_only=True)
                sheets_text = []
                for sheet in wb.worksheets:
                    sheet_data = []
                    sheet_data.append(f"Sheet: {sheet.title}")
                    # Simple extraction: iterate rows
                    for row in sheet.iter_rows(values_only=True):
                        # Filter out None values
                        row_data = [str(cell) for cell in row if cell is not None]
                        if row_data:
                            sheet_data.append("\t".join(row_data))
                    sheets_text.append("\n".join(sheet_data))
                text_content = "\n\n".join(sheets_text)

        elif file_ext == ".csv":
            # CSV needs text mode, but safe to read bytes and decode
            with file_field.open('r') as f:
                # If 'r' returns string (if backed by file), great. 
                # If bytes, we need to decode. 
                # Safe approach: read bytes, decode, then use io.StringIO
                pass
            
            with file_field.open('rb') as f:
                content_bytes = f.read()
                try:
                    content_str = content_bytes.decode('utf-8')
                except:
                    content_str = content_bytes.decode('latin-1', errors='replace')
                
                f_io = io.StringIO(content_str)
                reader = csv.reader(f_io)
                rows = []
                for row in reader:
                    rows.append(", ".join(row))
                text_content = "\n".join(rows)
        
        elif file_ext == ".txt" or file_ext in ['.py', '.js', '.html', '.css', '.java', '.c', '.cpp', '.h', '.md', '.json', '.xml', 'sql', '.sh', '.bat']:
             # Read as bytes and decode
             try:
                 with file_field.open('rb') as f:
                     content_bytes = f.read()
                     # Try utf-8 first
                     try:
                         text_content = content_bytes.decode('utf-8')
                     except UnicodeDecodeError:
                         text_content = content_bytes.decode('latin-1', errors='replace')
             except Exception as e:
                  print(f"Error reading text file {file_field.name}: {e}")
                  text_content = f"[Error reading file {file_field.name}]"


        else:
            # Fallback for other types: try to read as text, if fails, just mention file name
            try:
                with file_field.open('r') as f: # Try text mode first
                     text_content = f.read(10000) # Read first 10k chars
                     if len(text_content) == 10000:
                         text_content += "\n... (truncated)"
            except Exception:
                text_content = f"[File: {os.path.basename(filename)} (Type: {file_ext}) - Content not extractable]"

    except Exception as e:
        print(f"Error extracting file {file_field.name}: {e}")
        traceback.print_exc()
        text_content = f"[Error reading file {file_field.name}: {str(e)}]"

    return text_content, image_urls
